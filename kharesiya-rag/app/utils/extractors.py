import io
from dataclasses import dataclass

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class ExtractedPage:
    page_number: int
    text: str


def extract_text_from_pdf(content: bytes) -> list[ExtractedPage]:
    reader = PdfReader(io.BytesIO(content))
    pages: list[ExtractedPage] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append(ExtractedPage(page_number=i, text=text))
    return pages


def extract_text_from_docx(content: bytes) -> list[ExtractedPage]:
    doc = DocxDocument(io.BytesIO(content))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [ExtractedPage(page_number=1, text=text)]


def extract_text_from_pptx(content: bytes) -> list[ExtractedPage]:
    prs = Presentation(io.BytesIO(content))
    pages: list[ExtractedPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        pages.append(ExtractedPage(page_number=i, text="\n".join(texts)))
    return pages


def extract_text_from_plain(content: bytes) -> list[ExtractedPage]:
    text = content.decode("utf-8", errors="replace")
    return [ExtractedPage(page_number=1, text=text)]


def extract_text(content: bytes, extension: str) -> list[ExtractedPage]:
    ext = extension.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(content)
    if ext == ".docx":
        return extract_text_from_docx(content)
    if ext == ".pptx":
        return extract_text_from_pptx(content)
    if ext in {".txt", ".md", ".markdown"}:
        return extract_text_from_plain(content)
    raise ValueError(f"Unsupported extension: {extension}")
